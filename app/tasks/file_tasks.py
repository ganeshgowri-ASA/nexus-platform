"""
File Processing Tasks for Nexus Platform
Handles Word, Excel, PowerPoint, PDF, and image processing
"""
import os
import shutil
from pathlib import Path
from typing import Dict, Any, List, Optional
import logging
from datetime import datetime, timedelta

from celery import Task, group
from config.celery_config import celery_app
from config.settings import settings

logger = logging.getLogger(__name__)


class FileTask(Task):
    """Base class for file processing tasks"""
    autoretry_for = (IOError, OSError)
    retry_kwargs = {'max_retries': 2, 'countdown': 30}


@celery_app.task(base=FileTask, bind=True, name='app.tasks.file_tasks.process_word_document')
def process_word_document(self, file_path: str, operations: List[str]) -> Dict[str, Any]:
    """
    Process Word document (DOCX) with various operations

    Args:
        file_path: Path to the Word document
        operations: List of operations to perform (extract_text, convert_pdf, etc.)

    Returns:
        Processing results
    """
    try:
        from docx import Document

        doc = Document(file_path)
        results = {
            'file_path': file_path,
            'operations_completed': [],
            'outputs': {}
        }

        for operation in operations:
            if operation == 'extract_text':
                # Extract all text from document
                text = '\n'.join([para.text for para in doc.paragraphs])
                output_path = Path(file_path).with_suffix('.txt')
                output_path.write_text(text)
                results['outputs']['text_file'] = str(output_path)
                results['operations_completed'].append('extract_text')

            elif operation == 'count_stats':
                # Count paragraphs, tables, images
                stats = {
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'sections': len(doc.sections),
                    'total_chars': sum(len(p.text) for p in doc.paragraphs)
                }
                results['outputs']['stats'] = stats
                results['operations_completed'].append('count_stats')

            elif operation == 'extract_images':
                # Extract images from document
                image_dir = settings.TEMP_DIR / f"images_{Path(file_path).stem}"
                image_dir.mkdir(exist_ok=True)

                image_count = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        image_count += 1
                        image_data = rel.target_part.blob
                        image_path = image_dir / f"image_{image_count}.png"
                        image_path.write_bytes(image_data)

                results['outputs']['images_extracted'] = image_count
                results['outputs']['image_dir'] = str(image_dir)
                results['operations_completed'].append('extract_images')

        logger.info(f"Processed Word document: {file_path}")
        return results

    except Exception as e:
        logger.error(f"Error processing Word document: {str(e)}")
        raise self.retry(exc=e)


@celery_app.task(base=FileTask, bind=True, name='app.tasks.file_tasks.process_excel_file')
def process_excel_file(self, file_path: str, operations: List[str]) -> Dict[str, Any]:
    """
    Process Excel file (XLSX) with various operations

    Args:
        file_path: Path to the Excel file
        operations: List of operations to perform

    Returns:
        Processing results
    """
    try:
        from openpyxl import load_workbook
        import pandas as pd

        wb = load_workbook(file_path)
        results = {
            'file_path': file_path,
            'operations_completed': [],
            'outputs': {}
        }

        for operation in operations:
            if operation == 'analyze_sheets':
                # Analyze all sheets
                sheets_info = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheets_info.append({
                        'name': sheet_name,
                        'rows': sheet.max_row,
                        'columns': sheet.max_column
                    })
                results['outputs']['sheets'] = sheets_info
                results['operations_completed'].append('analyze_sheets')

            elif operation == 'convert_to_csv':
                # Convert each sheet to CSV
                csv_files = []
                for sheet_name in wb.sheetnames:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    csv_path = settings.TEMP_DIR / f"{Path(file_path).stem}_{sheet_name}.csv"
                    df.to_csv(csv_path, index=False)
                    csv_files.append(str(csv_path))
                results['outputs']['csv_files'] = csv_files
                results['operations_completed'].append('convert_to_csv')

            elif operation == 'calculate_totals':
                # Calculate totals for numeric columns
                totals = {}
                for sheet_name in wb.sheetnames:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    totals[sheet_name] = df[numeric_cols].sum().to_dict()
                results['outputs']['totals'] = totals
                results['operations_completed'].append('calculate_totals')

        logger.info(f"Processed Excel file: {file_path}")
        return results

    except Exception as e:
        logger.error(f"Error processing Excel file: {str(e)}")
        raise self.retry(exc=e)


@celery_app.task(base=FileTask, bind=True, name='app.tasks.file_tasks.process_powerpoint')
def process_powerpoint(self, file_path: str, operations: List[str]) -> Dict[str, Any]:
    """
    Process PowerPoint presentation (PPTX)

    Args:
        file_path: Path to the PowerPoint file
        operations: List of operations to perform

    Returns:
        Processing results
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        results = {
            'file_path': file_path,
            'operations_completed': [],
            'outputs': {}
        }

        for operation in operations:
            if operation == 'extract_text':
                # Extract text from all slides
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    slides_text.append({
                        'slide_number': i + 1,
                        'text': '\n'.join(slide_text)
                    })

                text_output = settings.TEMP_DIR / f"{Path(file_path).stem}_slides.txt"
                with open(text_output, 'w') as f:
                    for slide in slides_text:
                        f.write(f"\n--- Slide {slide['slide_number']} ---\n")
                        f.write(slide['text'])

                results['outputs']['slides_text'] = slides_text
                results['outputs']['text_file'] = str(text_output)
                results['operations_completed'].append('extract_text')

            elif operation == 'analyze_structure':
                # Analyze presentation structure
                structure = {
                    'total_slides': len(prs.slides),
                    'slide_layouts': {},
                    'total_shapes': 0
                }

                for slide in prs.slides:
                    layout_name = slide.slide_layout.name
                    structure['slide_layouts'][layout_name] = \
                        structure['slide_layouts'].get(layout_name, 0) + 1
                    structure['total_shapes'] += len(slide.shapes)

                results['outputs']['structure'] = structure
                results['operations_completed'].append('analyze_structure')

        logger.info(f"Processed PowerPoint file: {file_path}")
        return results

    except Exception as e:
        logger.error(f"Error processing PowerPoint file: {str(e)}")
        raise self.retry(exc=e)


@celery_app.task(base=FileTask, bind=True, name='app.tasks.file_tasks.process_pdf')
def process_pdf(self, file_path: str, operations: List[str]) -> Dict[str, Any]:
    """
    Process PDF file

    Args:
        file_path: Path to the PDF file
        operations: List of operations to perform

    Returns:
        Processing results
    """
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        results = {
            'file_path': file_path,
            'operations_completed': [],
            'outputs': {}
        }

        for operation in operations:
            if operation == 'extract_text':
                # Extract text from all pages
                text = []
                for i, page in enumerate(reader.pages):
                    text.append({
                        'page': i + 1,
                        'text': page.extract_text()
                    })

                text_output = settings.TEMP_DIR / f"{Path(file_path).stem}_extracted.txt"
                with open(text_output, 'w') as f:
                    for page in text:
                        f.write(f"\n--- Page {page['page']} ---\n")
                        f.write(page['text'])

                results['outputs']['text'] = text
                results['outputs']['text_file'] = str(text_output)
                results['operations_completed'].append('extract_text')

            elif operation == 'get_metadata':
                # Extract PDF metadata
                metadata = {
                    'pages': len(reader.pages),
                    'info': dict(reader.metadata) if reader.metadata else {}
                }
                results['outputs']['metadata'] = metadata
                results['operations_completed'].append('get_metadata')

        logger.info(f"Processed PDF file: {file_path}")
        return results

    except Exception as e:
        logger.error(f"Error processing PDF file: {str(e)}")
        raise self.retry(exc=e)


@celery_app.task(base=FileTask, bind=True, name='app.tasks.file_tasks.process_image')
def process_image(self, file_path: str, operations: List[str]) -> Dict[str, Any]:
    """
    Process image file

    Args:
        file_path: Path to the image file
        operations: List of operations to perform (resize, thumbnail, convert, etc.)

    Returns:
        Processing results
    """
    try:
        from PIL import Image

        img = Image.open(file_path)
        results = {
            'file_path': file_path,
            'operations_completed': [],
            'outputs': {}
        }

        for operation in operations:
            if operation == 'create_thumbnail':
                # Create thumbnail
                img_copy = img.copy()
                img_copy.thumbnail((200, 200))
                thumb_path = settings.TEMP_DIR / f"{Path(file_path).stem}_thumb.png"
                img_copy.save(thumb_path)
                results['outputs']['thumbnail'] = str(thumb_path)
                results['operations_completed'].append('create_thumbnail')

            elif operation == 'resize':
                # Resize to specific dimensions
                img_copy = img.copy()
                img_copy = img_copy.resize((800, 600))
                resized_path = settings.TEMP_DIR / f"{Path(file_path).stem}_resized.png"
                img_copy.save(resized_path)
                results['outputs']['resized'] = str(resized_path)
                results['operations_completed'].append('resize')

            elif operation == 'get_info':
                # Get image information
                info = {
                    'format': img.format,
                    'mode': img.mode,
                    'size': img.size,
                    'width': img.width,
                    'height': img.height
                }
                results['outputs']['info'] = info
                results['operations_completed'].append('get_info')

        logger.info(f"Processed image file: {file_path}")
        return results

    except Exception as e:
        logger.error(f"Error processing image file: {str(e)}")
        raise self.retry(exc=e)


@celery_app.task(name='app.tasks.file_tasks.batch_process_files')
def batch_process_files(files: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Process multiple files in parallel

    Args:
        files: List of dicts with 'path', 'type', and 'operations'

    Returns:
        Batch processing results
    """
    task_group = []

    for file_info in files:
        file_type = file_info['type']
        file_path = file_info['path']
        operations = file_info.get('operations', [])

        if file_type == 'word':
            task_group.append(process_word_document.s(file_path, operations))
        elif file_type == 'excel':
            task_group.append(process_excel_file.s(file_path, operations))
        elif file_type == 'powerpoint':
            task_group.append(process_powerpoint.s(file_path, operations))
        elif file_type == 'pdf':
            task_group.append(process_pdf.s(file_path, operations))
        elif file_type == 'image':
            task_group.append(process_image.s(file_path, operations))

    job = group(task_group)
    result = job.apply_async()

    return {
        'status': 'processing',
        'total_files': len(files),
        'group_id': result.id
    }


@celery_app.task(name='app.tasks.file_tasks.cleanup_temp_files')
def cleanup_temp_files(max_age_hours: int = 24) -> Dict[str, Any]:
    """
    Clean up temporary files older than max_age_hours

    Args:
        max_age_hours: Maximum age of files to keep in hours

    Returns:
        Cleanup results
    """
    try:
        cutoff_time = datetime.now() - timedelta(hours=max_age_hours)
        deleted_count = 0
        deleted_size = 0

        for temp_file in settings.TEMP_DIR.glob('**/*'):
            if temp_file.is_file():
                file_mtime = datetime.fromtimestamp(temp_file.stat().st_mtime)
                if file_mtime < cutoff_time:
                    file_size = temp_file.stat().st_size
                    temp_file.unlink()
                    deleted_count += 1
                    deleted_size += file_size

        logger.info(f"Cleaned up {deleted_count} temp files ({deleted_size} bytes)")
        return {
            'status': 'success',
            'deleted_files': deleted_count,
            'deleted_size_bytes': deleted_size
        }

    except Exception as e:
        logger.error(f"Error cleaning up temp files: {str(e)}")
        return {
            'status': 'error',
            'message': str(e)
        }
