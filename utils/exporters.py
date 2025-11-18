"""Export utilities for various file formats"""
from pathlib import Path
from typing import List, Dict, Any
from datetime import datetime
from config.settings import settings

def export_to_pdf(content: Any, output_path: str, document_type: str = "generic") -> str:
    """
    Export content to PDF format

    Args:
        content: Content to export (dict, list, or string)
        output_path: Path for the output file
        document_type: Type of document (presentation, email, note, report)

    Returns:
        Path to the exported file
    """
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    if document_type == "presentation":
        # Export slides as PDF
        for i, slide in enumerate(content.get('slides', []), 1):
            title = Paragraph(f"<b>Slide {i}: {slide.get('title', 'Untitled')}</b>", styles['Heading1'])
            story.append(title)
            story.append(Spacer(1, 0.2*inch))

            content_text = slide.get('content', '')
            if isinstance(content_text, list):
                content_text = '\n'.join(f"• {item}" for item in content_text)

            para = Paragraph(content_text, styles['Normal'])
            story.append(para)
            story.append(PageBreak())

    elif document_type == "email":
        # Export email as PDF
        title = Paragraph(f"<b>Subject: {content.get('subject', 'No Subject')}</b>", styles['Heading1'])
        story.append(title)
        story.append(Spacer(1, 0.2*inch))

        from_text = Paragraph(f"From: {content.get('sender', 'Unknown')}", styles['Normal'])
        story.append(from_text)

        to_text = Paragraph(f"To: {', '.join(content.get('recipients', []))}", styles['Normal'])
        story.append(to_text)

        story.append(Spacer(1, 0.3*inch))

        body = Paragraph(content.get('body', ''), styles['Normal'])
        story.append(body)

    elif document_type == "note":
        # Export note as PDF
        title = Paragraph(f"<b>{content.get('title', 'Untitled Note')}</b>", styles['Heading1'])
        story.append(title)
        story.append(Spacer(1, 0.2*inch))

        body = Paragraph(content.get('content', ''), styles['Normal'])
        story.append(body)

    else:
        # Generic export
        if isinstance(content, str):
            para = Paragraph(content, styles['Normal'])
            story.append(para)
        elif isinstance(content, dict):
            for key, value in content.items():
                heading = Paragraph(f"<b>{key}</b>", styles['Heading2'])
                story.append(heading)
                para = Paragraph(str(value), styles['Normal'])
                story.append(para)
                story.append(Spacer(1, 0.1*inch))

    doc.build(story)
    return str(output_path)

def export_to_pptx(slides: List[Dict], output_path: str, theme: str = "Modern Blue") -> str:
    """
    Export slides to PowerPoint format

    Args:
        slides: List of slide dictionaries with title and content
        output_path: Path for the output file
        theme: Theme name

    Returns:
        Path to the exported file
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        # Add a slide with title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Untitled')

        # Set content
        content = slide_data.get('content', '')
        if isinstance(content, list):
            content = '\n'.join(f"• {item}" for item in content)

        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.text = content

    prs.save(str(output_path))
    return str(output_path)

def export_to_xlsx(data: List[Dict], output_path: str, sheet_name: str = "Sheet1") -> str:
    """
    Export data to Excel format

    Args:
        data: List of dictionaries representing rows
        output_path: Path for the output file
        sheet_name: Name of the Excel sheet

    Returns:
        Path to the exported file
    """
    import pandas as pd

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    df = pd.DataFrame(data)
    df.to_excel(str(output_path), sheet_name=sheet_name, index=False)

    return str(output_path)

def export_to_docx(content: Dict, output_path: str) -> str:
    """
    Export content to Word document format

    Args:
        content: Dictionary with title and content
        output_path: Path for the output file

    Returns:
        Path to the exported file
    """
    from docx import Document

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()

    # Add title
    doc.add_heading(content.get('title', 'Untitled'), 0)

    # Add content
    content_text = content.get('content', '')
    if isinstance(content_text, list):
        for item in content_text:
            doc.add_paragraph(item, style='List Bullet')
    else:
        doc.add_paragraph(content_text)

    doc.save(str(output_path))
    return str(output_path)

def export_to_ics(events: List[Dict], output_path: str) -> str:
    """
    Export calendar events to ICS format

    Args:
        events: List of event dictionaries
        output_path: Path for the output file

    Returns:
        Path to the exported file
    """
    from icalendar import Calendar, Event
    import pytz

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    cal = Calendar()
    cal.add('prodid', '-//NEXUS Calendar//nexus-platform//')
    cal.add('version', '2.0')

    for event_data in events:
        event = Event()
        event.add('summary', event_data.get('title', 'Untitled Event'))
        event.add('description', event_data.get('description', ''))
        event.add('dtstart', event_data.get('start_time'))
        event.add('dtend', event_data.get('end_time'))
        event.add('location', event_data.get('location', ''))
        cal.add_component(event)

    with open(output_path, 'wb') as f:
        f.write(cal.to_ical())

    return str(output_path)
